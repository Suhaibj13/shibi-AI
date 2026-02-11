"""
GAIA V6 — File→AI Router (general-purpose, no audit tone, no forced citations)
- Text files:
    * If estimated tokens > THRESHOLD_TOKENS: compress with a cheaper pass, then answer from merged notes
    * Otherwise send raw (clamped) content
- Data files (csv/xlsx/tsv):
    * Plan DuckDB SQL → execute (if duckdb present) → summarize/polish answer

Style:
- "simple" (default): conversational/concise; NO citations/tables unless asked
- "structured": concise answer + brief bullet points (still no citations unless user asks)
"""

from __future__ import annotations
import io, os, json, hashlib, mimetypes
from dataclasses import dataclass
from typing import List, Tuple, Optional, Dict

# Optional parsers (best-effort)
try:
    import pandas as pd
except Exception:
    pd = None

try:
    import duckdb  # for executing SQL plans
except Exception:
    duckdb = None

# --- Model resolution ---
from models import resolve, ResolvedModel

# ===== Config =====
THRESHOLD_TOKENS = 1000           # if estimated tokens > threshold → token-efficient route
MAX_DIRECT_INPUT_TOKENS = 2000    # cap when sending raw content
MAX_RESULT_ROWS = 200             # cap table rows from SQL engine
ROW_SAMPLE_FOR_SUMMARY = 60       # if > this, cheap model summarizes table first

TEXT_EXT = {".txt", ".md", ".pdf", ".docx", ".pptx", ".rtf"}
DATA_EXT = {".csv", ".tsv", ".xlsx", ".xls"}

# If you want a specific cheaper model per provider, set here
CHEAP_FALLBACK_BY_PROVIDER = {
    "groq":      "llama-3.1-8b-instant",
    "openai":    "gpt-4o",
    "anthropic": "claude-3-haiku-20240307",
    "gemini":    "gemini-2.5-flash",
    "cohere":    "command-r",
}


# ===== Utilities =====
def sha256_bytes(b: bytes) -> str:
    h = hashlib.sha256(); h.update(b); return h.hexdigest()

def ext_of(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower()

def estimate_tokens_from_text(text: str) -> int:
    # very rough: ~4 chars/token
    return max(1, int(len(text) / 4))

def clamp_tokens(text: str, max_tokens: int) -> str:
    max_chars = max_tokens * 4
    return text if len(text) <= max_chars else text[:max_chars]

# ===== File wrappers =====
@dataclass
class UploadedFile:
    name: str
    bytes: bytes
    mime: str
    sha256: str
    ext: str

def wrap_uploaded(file_storage) -> UploadedFile:
    data = file_storage.read()
    file_storage.stream.seek(0)
    name = getattr(file_storage, "filename", "upload.bin") or "upload.bin"
    mime = file_storage.mimetype or (mimetypes.guess_type(name)[0] or "application/octet-stream")
    return UploadedFile(name=name, bytes=data, mime=mime, sha256=sha256_bytes(data), ext=ext_of(name))

def is_text_file(f: UploadedFile) -> bool:
    if f.ext in TEXT_EXT: return True
    if f.mime.startswith("text/"): return True
    return False

def is_data_file(f: UploadedFile) -> bool:
    if f.ext in DATA_EXT: return True
    if f.mime in ("text/csv", "text/tab-separated-values", "application/vnd.ms-excel",
                  "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"):
        return True
    return False

# ===== Text extraction (best-effort) =====
def extract_text_best_effort(f: UploadedFile) -> str:
    try:
        if f.ext == ".txt" or f.mime.startswith("text/"):
            return f.bytes.decode("utf-8", errors="replace")

        if f.ext == ".md":
            return f.bytes.decode("utf-8", errors="replace")

        if f.ext == ".pdf":
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(f.bytes))
                parts = []
                for i, p in enumerate(reader.pages or []):
                    parts.append(f"[page {i+1}]\n" + (p.extract_text() or ""))
                return "\n\n".join(parts).strip()
            except Exception:
                return ""

        if f.ext == ".docx":
            try:
                import docx
                d = docx.Document(io.BytesIO(f.bytes))
                return "\n".join([p.text for p in d.paragraphs])
            except Exception:
                return ""

        if f.ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(f.bytes))
                slides = []
                for i, s in enumerate(prs.slides):
                    texts = []
                    for shape in s.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                    slides.append(f"[slide {i+1}]\n" + "\n".join(texts))
                return "\n\n".join(slides)
            except Exception:
                return ""
    except Exception:
        pass
    return ""

# ===== Data loading (pandas/duckdb) =====
@dataclass
class TableProfile:
    table_name: str
    columns: List[Tuple[str, str]]  # (name, dtype)
    quick_stats: Dict[str, str]     # col -> brief
    row_count: int
    engine_table_ref: str           # e.g., 't0'

def load_table(f: UploadedFile, duck_conn=None) -> Optional[TableProfile]:
    if pd is None:
        return None
    try:
        if f.ext in (".csv", ".tsv"):
            sep = "," if f.ext == ".csv" else "\t"
            df = pd.read_csv(io.BytesIO(f.bytes), sep=sep, nrows=200000)
        elif f.ext in (".xlsx", ".xls"):
            df = pd.read_excel(io.BytesIO(f.bytes))
        else:
            return None
    except Exception:
        return None

    cols = [(c, str(t)) for c, t in zip(df.columns, df.dtypes)]
    quick = {}
    for c in df.columns[:25]:
        try:
            nonnull = df[c].notna().sum()
            nullpct = 0.0 if len(df) == 0 else round((len(df) - nonnull) * 100.0 / len(df), 1)
            sample = df[c].dropna().astype(str).head(3).tolist()
            quick[c] = f"nonnull%={100-nullpct} samples={sample}"
        except Exception:
            quick[c] = "n/a"

    table_ref = "t0"
    if duckdb is not None and duck_conn is not None:
        try:
            duck_conn.register(table_ref, df)
        except Exception:
            table_ref = "t0"

    return TableProfile(table_name=f.name, columns=cols, quick_stats=quick, row_count=len(df), engine_table_ref=table_ref)

# ===== Model calls (centralized) =====
def choose_cheap_concrete(rm: ResolvedModel) -> Tuple[str, str]:
    cheap_model = CHEAP_FALLBACK_BY_PROVIDER.get(rm.provider)
    return (rm.provider, cheap_model or rm.model_id)

def call_chat_completion(provider: str, model_id: str, system: str, user: str, temperature: float = 0.2) -> str:
    if provider == "groq":
        from groq import Groq
        client = Groq(api_key=os.getenv("GROQ_API_KEY"))
        resp = client.chat.completions.create(
            model=model_id,
            temperature=temperature,
            messages=[{"role": "system", "content": system}, {"role": "user", "content": user}]
        )
        return (resp.choices[0].message.content or "").strip()

    if provider == "openai":
        api_key = os.getenv("GPT_API_KEY") or os.getenv("OPENAI_API_KEY") or os.getenv("OPENAI_KEY")
        if not api_key:
            raise RuntimeError("OpenAI API key missing. Set GPT_API_KEY or OPENAI_API_KEY.")

        # detect SDK version
        try:
            import openai
            ver = getattr(openai, "__version__", "0")
            major = int(str(ver).split(".", 1)[0])
        except Exception:
            major = 0

        if major >= 1:
            try:
                from openai import OpenAI
                client = OpenAI(api_key=api_key)
            except Exception:
                import openai as openai_mod
                client = openai_mod.OpenAI(api_key=api_key)

            resp = client.chat.completions.create(
                model=model_id,
                temperature=temperature,
                messages=[{"role":"system","content":system},{"role":"user","content":user}]
            )
            return (resp.choices[0].message.content or "").strip()

        # legacy (<1.0)
        import openai as openai_legacy
        openai_legacy.api_key = api_key
        r = openai_legacy.ChatCompletion.create(
            model=model_id,
            messages=[{"role":"system","content":system},{"role":"user","content":user}],
            temperature=temperature
        )
        return (r["choices"][0]["message"]["content"] or "").strip()


    if provider == "anthropic":
        try:
            import anthropic
        except Exception:
            raise RuntimeError("anthropic package not installed.")
        client = anthropic.Anthropic(api_key=os.getenv("CLAUDE_API_KEY"))
        r = client.messages.create(
            model=model_id, temperature=temperature, max_tokens=1024,
            messages=[
                {"role":"user","content":[{"type":"text","text": f"SYSTEM:\n{system}\n\nUSER:\n{user}"}]}
            ]
        )
        parts = []
        for b in r.content or []:
            t = getattr(b, "text", None) or (b.get("text") if isinstance(b, dict) else None)
            if t: parts.append(t)
        return "\n".join(parts).strip()

    if provider == "gemini":
        try:
            import google.generativeai as genai
        except Exception:
            raise RuntimeError("google-generativeai package not installed.")
        genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
        model = genai.GenerativeModel(model_id)
        chat = model.start_chat(history=[])
        resp = chat.send_message(f"SYSTEM:\n{system}\n\nUSER:\n{user}")
        return (getattr(resp, "text", "") or "").strip()

    if provider == "cohere":
        try:
            from cohere import ClientV2
            c2 = ClientV2(api_key=os.getenv("COHERE_API_KEY"))
            r = c2.responses.create(model=model_id, input=f"SYSTEM:\n{system}\n\nUSER:\n{user}")
            return (getattr(r, "output_text", None) or getattr(r, "text", "") or "").strip()
        except Exception:
            import cohere
            client = cohere.Client(os.getenv("COHERE_API_KEY"))
            r = client.chat(model=model_id, message=f"SYSTEM:\n{system}\n\nUSER:\n{user}")
            return (getattr(r, "text", None) or getattr(r, "generation", "") or "").strip()

    raise RuntimeError(f"Provider '{provider}' not implemented.")


# ===== Compression prompts (general) =====
def cheap_compress_chunk(rm: ResolvedModel, file_name: str, ref: str, chunk_text: str) -> str:
    prov, cheap_model = choose_cheap_concrete(rm)
    system = (
        "You compress text into clear, factual bullet notes. "
        "Preserve names, dates, amounts, and units. Do not invent or add citations."
    )
    user = f"File: {file_name}  Ref: {ref}\nText:\n<<<{chunk_text}>>>\nCompress to ≤120 tokens of bullets."
    return call_chat_completion(prov, cheap_model, system, user, temperature=0.0)

def cheap_merge_notes(rm: ResolvedModel, notes: List[str], file_name: str) -> str:
    prov, cheap_model = choose_cheap_concrete(rm)
    joined = "\n".join([n for n in notes if n])
    system = "Merge bullet notes. Dedupe and keep neutral tone. ≤800 tokens."
    user = f"File: {file_name}\nNotes:\n<<<{joined}>>>\nReturn merged notes ≤800 tokens, grouped by obvious sections."
    return call_chat_completion(prov, cheap_model, system, user, temperature=0.0)

# ===== Final answer (general styles) =====
def expensive_final_answer(
    rm: ResolvedModel,
    user_query: str,
    file_manifest: List[str],
    evidence_notes: str,
    style: str = "simple"
) -> str:
    style = (style or "simple").strip().lower()
    manifest = "\n".join(f"- {n}" for n in file_manifest)

    if style == "structured":
        system = (
            "You are GAIA, a helpful general-purpose assistant. "
            "Write concise, accurate answers. Avoid citations and tables unless explicitly requested."
        )
        user = f"""Question: {user_query}
Files:
{manifest}

Notes from files:
<<<
{evidence_notes}
>>>

Instructions:
- Start with a 2–4 line direct answer.
- Then 3–5 brief bullets with key points or caveats.
- If information is missing, say so plainly.
"""
        return call_chat_completion(rm.provider, rm.model_id, system, user, temperature=0.2)

    # default: simple
    system = (
        "You are GAIA, a helpful general-purpose assistant. "
        "Respond conversationally and concisely. Do NOT add citations or tables unless asked."
    )
    user = f"""Question: {user_query}
Files:
{manifest}

Notes from files:
<<<
{evidence_notes}
>>>

Instructions:
- 3–6 lines, natural tone.
- No citations. No tables.
- If something is unclear or missing, state it briefly."""
    return call_chat_completion(rm.provider, rm.model_id, system, user, temperature=0.2)

# ===== Data path prompts (general) =====
def expensive_plan_sql(rm: ResolvedModel, user_query: str, profile: TableProfile) -> Dict:
    system = "You generate precise DuckDB SQL plans as VALID JSON only."
    cols_str = ", ".join([f"{c}:{t}" for c, t in profile.columns])
    stats_lines = "\n".join([f"- {k}: {v}" for k, v in profile.quick_stats.items()])
    user = f"""Goal: {user_query}
Table: {profile.table_name}
Schema: {cols_str}
Quick stats:
{stats_lines}

Constraints:
- Prefer DuckDB SQL with CTEs if useful.
- Output JSON ONLY:
{{
  "rationale": "...",
  "sql": "...",
  "expects_rows": true/false,
  "columns_used": [...]
}}"""
    raw = call_chat_completion(rm.provider, rm.model_id, system, user, temperature=0.1)
    try:
        return json.loads(raw)
    except Exception:
        return {"rationale": "Failed to parse plan JSON.", "sql": "", "expects_rows": True, "columns_used": []}

def cheap_summarize_rows(rm: ResolvedModel, user_query: str, table_name: str, rows_preview: str) -> str:
    prov, cheap_model = choose_cheap_concrete(rm)
    system = "Summarize the tabular results into ≤150 tokens. Keep exact figures that appear. No citations."
    user = f"Question: {user_query}\nResult preview from {table_name}:\n<<<{rows_preview}>>>\nReturn a concise factual summary."
    return call_chat_completion(prov, cheap_model, system, user, temperature=0.0)

def expensive_polish_result(rm: ResolvedModel, user_query: str, plan_json: Dict, result_block: str, style: str = "simple") -> str:
    style = (style or "simple").strip().lower()
    if style == "structured":
        system = "You are GAIA. Write a crisp, decision-ready answer based on query, plan, and result. No citations."
        user = f"""Question: {user_query}
Plan (JSON):
{json.dumps(plan_json, indent=2)}
Result:
<<<{result_block}>>>

Instructions:
- Start with the answer in 2–4 lines.
- Then 3–5 key takeaways (bullets).
- If sampling/caps applied, mention briefly."""
        return call_chat_completion(rm.provider, rm.model_id, system, user, temperature=0.2)

    # simple
    system = "You are GAIA. Write a short, friendly answer from the result. No citations."
    user = f"""Question: {user_query}
Result:
<<<{result_block}>>>

Instructions:
- 3–6 lines, natural tone.
- No citations/tables unless requested."""
    return call_chat_completion(rm.provider, rm.model_id, system, user, temperature=0.2)

# ===== Chunking =====
def make_chunks(text: str, chunk_chars: int = 4000, overlap: int = 400) -> List[Tuple[str, str]]:
    chunks = []
    n = len(text)
    if n == 0: return chunks
    i = 0; idx = 1
    while i < n:
        j = min(n, i + chunk_chars)
        chunk = text[i:j]
        ref = f"part-{idx}"
        chunks.append((ref, chunk))
        i = max(i + chunk_chars - overlap, j)
        idx += 1
    return chunks

# ===== Main router =====
def route_and_answer(
    user_query: str,
    files: List,                 # list of Flask file storages (request.files.getlist("files[]"))
    logical_model_key: str = "grok",
    style: str = "simple"        # "simple" | "structured"
) -> Dict:
    """
    Entry point. Returns { "reply": str, "meta": {...} }
    """
    style = (style or "simple").strip().lower()
    rm = resolve(logical_model_key)
    uploaded: List[UploadedFile] = [wrap_uploaded(f) for f in (files or [])]

    if not uploaded:
        return {"reply": "No files received. Please attach a file.", "meta": {"files": 0}}

    text_files = [f for f in uploaded if is_text_file(f)]
    data_files = [f for f in uploaded if is_data_file(f)]

    manifest_names = [f.name for f in uploaded]
    meta_out = {"files": len(uploaded), "model": rm.logical, "manifest": manifest_names}

    final_text_evidence: List[str] = []
    final_answers: List[str] = []

    # ---- TEXT PATH ----
    for f in text_files:
        text = extract_text_best_effort(f)
        if not text:
            continue
        est = estimate_tokens_from_text(text)
        if est > THRESHOLD_TOKENS:
            notes = []
            for ref, chunk in make_chunks(text):
                notes.append(cheap_compress_chunk(rm, f.name, ref, chunk))
            merged = cheap_merge_notes(rm, notes, f.name)
            final_text_evidence.append(merged)
        else:
            final_text_evidence.append(f"[{f.name}]\n" + clamp_tokens(text, MAX_DIRECT_INPUT_TOKENS))

    # ---- DATA PATH ----
    duck_conn = duckdb.connect() if duckdb is not None else None
    try:
        for f in data_files:
            profile = load_table(f, duck_conn)
            if profile is None:
                continue
            plan = expensive_plan_sql(rm, user_query, profile)
            sql = (plan.get("sql") or "").strip()
            if not sql or duck_conn is None:
                final_answers.append(expensive_polish_result(rm, user_query, plan, "Execution unavailable.", style=style))
                continue

            try:
                rel = duck_conn.sql(sql)
                df = rel.df()
                if len(df.columns) > 12:
                    df = df.iloc[:, :12]
                if len(df) > MAX_RESULT_ROWS:
                    df_preview = df.head(MAX_RESULT_ROWS)
                    too_big = True
                else:
                    df_preview = df
                    too_big = False

                if too_big or len(df_preview) > ROW_SAMPLE_FOR_SUMMARY:
                    rows_text = df_preview.to_csv(index=False)
                    summary = cheap_summarize_rows(rm, user_query, profile.table_name, rows_text)
                    result_block = summary
                else:
                    result_block = df_preview.to_csv(index=False)

                polished = expensive_polish_result(rm, user_query, plan, result_block, style=style)
                final_answers.append(polished)

            except Exception as e:
                final_answers.append(expensive_polish_result(
                    rm, user_query, plan, f"SQL execution failed: {e}", style=style
                ))
    finally:
        try:
            if duck_conn is not None:
                duck_conn.close()
        except Exception:
            pass

    # ---- Final assembly ----
    if final_answers:
        # NEW
        return {"reply": "\n\n".join(final_answers), "meta": meta_out, "model_used": rm.model_id}

    evidence_notes = "\n\n".join(final_text_evidence).strip()
    if evidence_notes:
        reply = expensive_final_answer(rm, user_query, manifest_names, evidence_notes, style=style)
        # NEW
        return {"reply": reply, "meta": meta_out, "model_used": rm.model_id}


    # NEW
    return {"reply": "I couldn't read useful content from the attached files. Please try another format.",
        "meta": meta_out, "model_used": rm.model_id}
